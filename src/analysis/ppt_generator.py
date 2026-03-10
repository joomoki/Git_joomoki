import sys
import os
import io
import pandas as pd
import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import mplfinance as mpf
import matplotlib.pyplot as plt

# 프로젝트 루트 경로 추가
sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from config.db_config import SCHEMA_NAME
from src.stock_db_manager import StockDBManager
from src.analysis.indicators import analyze_stock

def create_ppt_report():
    print("=== 스마트폰용 주식 분석 PPT 생성 시작 ===")
    
    db_manager = StockDBManager()
    if not db_manager.connect():
        print("DB 연결 실패")
        return

    try:
        # 1. 매수 추천 종목 조회 (일단 최근 분석 결과나 직접 분석)
        # 상위 10개만 추출 (점수 높은 순)
        # 여기서 get_market_overview 재활용하되, 직접 데이터 긁어와서 다시 분석할 수도 있음.
        # 가장 확실한 건 DB에서 코드를 가져와서 다시 분석하는 것 (지표가 업데이트 되었으므로)
        
        # 전체 종목 중 100개만 샘플링하거나, 이미 분석된 것 중 'UP'인 것만 가져옴
        with db_manager.conn.cursor() as cur:
             # 예측이 UP인 종목 중 시가총액 상위 50개 가져오기 (우량주 위주 추천)
             sql = f"""
                SELECT 
                    c.stock_code, c.company_name, c.market_type
                FROM {SCHEMA_NAME}.stock_companies c
                JOIN {SCHEMA_NAME}.stock_analysis a ON c.stock_code = a.stock_code
                WHERE a.price_prediction = 'UP' 
                AND a.analysis_date = (SELECT MAX(analysis_date) FROM {SCHEMA_NAME}.stock_analysis)
                ORDER BY c.market_cap DESC
                LIMIT 10
             """
             cur.execute(sql)
             target_stocks = cur.fetchall()
             
        if not target_stocks:
            print("매수 추천 종목이 없습니다. (분석 데이터가 없거나 매수 신호 없음)")
            # 임시로 랜덤 3개라도 가져오기
            with db_manager.conn.cursor() as cur:
                 cur.execute(f"SELECT stock_code, company_name, market_type FROM {SCHEMA_NAME}.stock_companies ORDER BY market_cap DESC LIMIT 5")
                 target_stocks = cur.fetchall()

        # PPT 생성
        prs = Presentation()
        
        # 스마트폰 비율 (16:9 세로? 아니면 그냥 16:9 가로가 나음. 폰은 돌려볼 수 있으니 가로 추천)
        # 기본 16:9 사용.
        
        # 타이틀 슬라이드
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "주모키 AI 추천 종목 Top 10"
        subtitle.text = f"기준일: {datetime.datetime.now().strftime('%Y-%m-%d')}\n스마트폰 최적화 리포트"

        for stock_code, stock_name, market_type in target_stocks:
            print(f"분석 중: {stock_name} ({stock_code})")
            
            # 데이터 가져오기 (100일)
            with db_manager.conn.cursor() as cur:
                cur.execute(f"""
                    SELECT trade_date, open_price, high_price, low_price, close_price, volume 
                    FROM {SCHEMA_NAME}.stock_prices 
                    WHERE stock_code = %s 
                    ORDER BY trade_date DESC 
                    LIMIT 100
                """, (stock_code,))
                rows = cur.fetchall()
                
            if not rows: continue
            
            rows.reverse()
            df = pd.DataFrame(rows, columns=['trade_date', 'open_price', 'high_price', 'low_price', 'close_price', 'volume'])
            df.index = pd.DatetimeIndex(df['trade_date'])

            # Decimal -> float 변환 (pandas-ta 오류 방지)
            numeric_cols = ['open_price', 'high_price', 'low_price', 'close_price', 'volume']
            df[numeric_cols] = df[numeric_cols].astype(float)
            
            # 분석 실행
            analysis = analyze_stock(df)
            
            # 차트 이미지 생성
            chart_filename = f"chart_{stock_code}.png"
            generate_chart_image(df, stock_name, chart_filename)
            
            # 슬라이드 추가 (빈 슬라이드에 커스텀 레이아웃)
            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 1. 제목 (종목명 + 가격) - 아주 크게
            left = Inches(0.5)
            top = Inches(0.2)
            width = Inches(9)
            height = Inches(1.5)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = f"{stock_name}"
            p.font.bold = True
            p.font.size = Pt(40)
            p.font.color.rgb = RGBColor(0, 0, 0)
            
            p2 = tf.add_paragraph()
            p2.text = f"{analysis['current_price']:,}원 ({stock_code})"
            p2.font.size = Pt(24)
            p2.font.color.rgb = RGBColor(80, 80, 80)

            # 2. 종합 의견 배지
            summary_color = RGBColor(255, 0, 0) if analysis['score'] >= 1 else RGBColor(0, 0, 255)
            
            left = Inches(7)
            top = Inches(0.4)
            width = Inches(2.5)
            height = Inches(0.8)
            shape = slide.shapes.add_shape(
                1, left, top, width, height # 1 = Rectangle
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = summary_color
            shape.line.color.rgb = summary_color
            
            text_frame = shape.text_frame
            text_frame.text = analysis['summary']
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.size = Pt(20)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # 3. 차트 이미지 삽입
            img_path = chart_filename
            if os.path.exists(img_path):
                left = Inches(0.5)
                top = Inches(2.0)
                width = Inches(6)
                height = Inches(4)
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)
                os.remove(img_path) # 사용 후 삭제

            # 4. 핵심 분석 포인트 (오른쪽 or 아래)
            left = Inches(6.8)
            top = Inches(2.0)
            width = Inches(3.0)
            height = Inches(4)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            
            # 헤더
            p = tf.paragraphs[0]
            p.text = "핵심 포인트"
            p.font.bold = True
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(50, 50, 50)
            
            # 신호들
            for signal in analysis['signals']:
                p = tf.add_paragraph()
                p.text = f"• {signal['msg']}"
                p.font.size = Pt(14)
                p.space_after = Pt(10)
                
            # 기본 데이터
            p = tf.add_paragraph()
            p.text = f"\nRSI: {analysis['rsi']}"
            p.font.size = Pt(14)
            
            p = tf.add_paragraph()
            p.text = f"거래량: {analysis['current_price'] * df['volume'].iloc[-1] // 100000000}억원"
            p.font.size = Pt(14)

        # 저장
        output_path = os.path.join("src", "web", "static", "stock_report.pptx")
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        print(f"PPT 생성 완료: {output_path}")

    except Exception as e:
        print(f"PPT 생성 중 오류: {e}")
    finally:
        db_manager.disconnect()

def generate_chart_image(df, title, filename):
    try:
        # mplfinance를 위한 컬럼명 변경
        chart_df = df.copy()
        chart_df = chart_df.rename(columns={
            'open_price': 'Open',
            'high_price': 'High',
            'low_price': 'Low',
            'close_price': 'Close',
            'volume': 'Volume'
        })
        
        # 이동평균선 등도 차트 데이터프레임에 있어야 함 (analyze_stock에서 추가됨)
        # 하지만 analyze_stock은 원본 df에 추가했으므로, chart_df에도 복사 필요
        # 아니면 위에서 복사했으니 이미 포함됨.
        
        # mplfinance 스타일 설정
        s = mpf.make_mpf_style(base_mpf_style='charles', rc={'font.size': 10})
        
        # 차트 그리기 (캔들 + 이동평균선 + 거래량)
        apds = []
        if 'MA5' in chart_df.columns:
            apds.append(mpf.make_addplot(chart_df['MA5'], color='green', width=1))
        if 'MA20' in chart_df.columns:
            apds.append(mpf.make_addplot(chart_df['MA20'], color='red', width=1))
        
        if 'BBU_20_2.0' in chart_df.columns:
             apds.append(mpf.make_addplot(chart_df['BBU_20_2.0'], color='gray', linestyle='dotted'))
             apds.append(mpf.make_addplot(chart_df['BBL_20_2.0'], color='gray', linestyle='dotted'))

        mpf.plot(
            chart_df, 
            type='candle', 
            style=s, 
            title='', # 한글 깨짐 방지 (슬라이드에 텍스트로 표시됨)
            ylabel='',
            volume=True, 
            addplot=apds,
            savefig=dict(fname=filename, dpi=100, pad_inches=0.2),
            figratio=(10, 6),
            figscale=1.0,
            tight_layout=True
        )
    except Exception as e:
        print(f"차트 생성 실패: {e}")

if __name__ == "__main__":
    create_ppt_report()
